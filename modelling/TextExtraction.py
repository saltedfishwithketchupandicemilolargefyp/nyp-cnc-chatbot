# import required libraries for handling different document types
from pdfminer.high_level import extract_text
import os
from docx import Document
from pptx import Presentation
import openpyxl

# function to extract text from pdf files
# uses pdfminer.six library to process each pdf and combine their text
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        text += extract_text(pdf)
    return text

# function to extract text from word documents (.docx)
# processes paragraphs, tables and embedded content
def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        
        # get text from regular paragraphs
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        # get text from tables in the document
        # processes cell by cell, row by row
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"  # add line break after each row
    
    return text

# function to extract text from powerpoint presentations (.pptx)
# handles slides, shapes, tables and speaker notes
def get_pptx_text(pptx_docs):
    text = ""
    
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        
        # process each slide in the presentation
        for slide in presentation.slides:
            
            # get text from all shapes (textboxes, titles, etc)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                
                # handle tables within slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + "\t"  # separate columns with tabs
                        text += "\n"  # new line for each row
            
            # extract speaker notes if they exist
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide:
                    notes_text = notes_slide.notes_text_frame.text
                    text += f"{notes_text}\n"
                
    return text

# function to extract text from excel files (.xlsx)
# processes all sheets and cells with content
def get_xlsx_text(xlsx_docs):
    text = ""
    for xlsx in xlsx_docs:
        wb = openpyxl.load_workbook(xlsx)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text += str(cell.value) + " "
                text += "\n"
    return text

# set up directory paths for different document types
pdf_dir = r'modelling\data\pdf_files'
docx_dir = r'modelling\data\docx_files'
pptx_dir = r'modelling\data\pptx_files'
xlsx_dir = r'modelling\data\xlsx_files'

# get list of files for each document type
# uses list comprehension to filter files by extension
pdf_files = [os.path.join(pdf_dir, f) for f in os.listdir(pdf_dir) if f.endswith('.pdf')]
docx_files = [os.path.join(docx_dir, f) for f in os.listdir(docx_dir) if f.endswith('.docx')]
pptx_files = [os.path.join(pptx_dir, f) for f in os.listdir(pptx_dir) if f.endswith('.pptx')]
xlsx_files = [os.path.join(xlsx_dir, f) for f in os.listdir(xlsx_dir) if f.endswith('.xlsx')]

# extract text from all document types
pdf_text = get_pdf_text(pdf_files)
docx_text = get_docx_text(docx_files)
pptx_text = get_pptx_text(pptx_files)
xlsx_text = get_xlsx_text(xlsx_files)

# set output file path for combined text
output_file = os.path.join('modelling', 'extracted_text.txt')

# function to save all extracted text to a single file
def save_all_text_to_file(texts, file_name):
    with open(file_name, 'w', encoding='utf-8') as f:
        for text_type, text in texts.items():
            # write text from each document type
            f.write(text + "\n")

# create dictionary of all extracted texts
all_texts = {
    "pdf": pdf_text,
    "docx": docx_text,
    "pptx": pptx_text,
    "xlsx": xlsx_text
}

# save all extracted text to the output file
save_all_text_to_file(all_texts, output_file)